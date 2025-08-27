from pptx import Presentation


def build_deck(df, out_path, max_results: int = 50):
    """Build a PowerPoint deck of the top funding opportunities.

    Parameters
    ----------
    df : pandas.DataFrame
        Scored funding opportunities. Expected to contain ``Rank``, ``Grant Name``
        and ``Weighted Score`` columns.
    out_path : str
        Path where the deck should be saved.
    max_results : int, optional
        Maximum number of opportunities to include. Defaults to 50.
    """

    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    slide.shapes.title.text = "EQORE Funding Deck"
    slide.placeholders[1].text = "Auto-generated deck"

    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "Top Opportunities"

    top_df = df.head(max_results)
    text = "\n".join(
        [
            f"{row['Rank']}. {row['Grant Name']} ({row['Weighted Score']})"
            for _, row in top_df.iterrows()
        ]
    )
    slide.placeholders[1].text = text

    prs.save(out_path)
